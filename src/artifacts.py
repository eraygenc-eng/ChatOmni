from pathlib import Path, PurePosixPath
import html
import json
import re
import zipfile

import reportlab
from docx import Document
from openpyxl import Workbook
from openpyxl.styles import Alignment, Font
from openpyxl.utils import get_column_letter
from pptx import Presentation
from pptx.util import Inches
from reportlab.lib import colors
from reportlab.lib.enums import TA_CENTER
from reportlab.lib.pagesizes import A4
from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
from reportlab.lib.units import mm
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.ttfonts import TTFont
from reportlab.platypus import Paragraph, SimpleDocTemplate, Spacer, Table, TableStyle


MAX_ZIP_FILES = 2000
MAX_ZIP_TEXT_BYTES = 25 * 1024 * 1024


def _text(value) -> str:
    if value is None:
        return ""
    if isinstance(value, str):
        return value
    if isinstance(value, (dict, list, tuple)):
        return json.dumps(value, ensure_ascii=False)
    return str(value)


def _table_parts(table_spec: dict):
    if not isinstance(table_spec, dict):
        raise ValueError("Each table must be an object.")

    headers = table_spec.get("headers") or table_spec.get("columns") or []
    rows = table_spec.get("rows") or []

    if not isinstance(headers, list) or not isinstance(rows, list):
        raise ValueError("Table headers and rows must be lists.")

    headers = [_text(value) for value in headers]
    normalized_rows = []

    for row in rows:
        if isinstance(row, dict):
            if headers:
                normalized_rows.append([row.get(header, "") for header in headers])
            else:
                normalized_rows.append(list(row.values()))
        elif isinstance(row, (list, tuple)):
            normalized_rows.append(list(row))
        else:
            normalized_rows.append([row])

    return headers, normalized_rows


def _document_blocks(spec: dict):
    root = {
        "paragraphs": spec.get("paragraphs") or [],
        "bullets": spec.get("bullets") or [],
        "numbered": spec.get("numbered") or [],
        "tables": spec.get("tables") or [],
    }

    sections = spec.get("sections") or []
    if not isinstance(sections, list):
        raise ValueError("'sections' must be a list.")

    return [root, *sections]


def _write_docx(path: Path, spec: dict):
    document = Document()

    title = _text(spec.get("title")).strip()
    subtitle = _text(spec.get("subtitle")).strip()

    if title:
        document.add_heading(title, level=0)

    if subtitle:
        paragraph = document.add_paragraph()
        run = paragraph.add_run(subtitle)
        run.italic = True

    content = _text(spec.get("content")).strip()
    if content:
        for paragraph_text in re.split(r"\n\s*\n", content):
            if paragraph_text.strip():
                document.add_paragraph(paragraph_text.strip())

    for block in _document_blocks(spec):
        if not isinstance(block, dict):
            raise ValueError("Each section must be an object.")

        heading = _text(block.get("heading") or block.get("title")).strip()
        if heading:
            document.add_heading(heading, level=1)

        paragraphs = block.get("paragraphs") or []
        if isinstance(paragraphs, str):
            paragraphs = [paragraphs]
        for item in paragraphs:
            document.add_paragraph(_text(item))

        bullets = block.get("bullets") or []
        if isinstance(bullets, str):
            bullets = [bullets]
        for item in bullets:
            document.add_paragraph(_text(item), style="List Bullet")

        numbered = block.get("numbered") or []
        if isinstance(numbered, str):
            numbered = [numbered]
        for item in numbered:
            document.add_paragraph(_text(item), style="List Number")

        tables = block.get("tables") or []
        if not isinstance(tables, list):
            raise ValueError("Section tables must be a list.")

        for table_spec in tables:
            headers, rows = _table_parts(table_spec)
            column_count = max([len(headers), *(len(row) for row in rows), 1])
            table = document.add_table(rows=1 if headers else 0, cols=column_count)
            table.style = "Table Grid"

            if headers:
                for index, header in enumerate(headers):
                    cell = table.rows[0].cells[index]
                    cell.text = header
                    for run in cell.paragraphs[0].runs:
                        run.bold = True

            for row in rows:
                cells = table.add_row().cells
                for index in range(column_count):
                    cells[index].text = _text(row[index] if index < len(row) else "")

    document.save(path)


def _pdf_fonts():
    candidates = [
        (
            Path(reportlab.__file__).resolve().parent / "fonts" / "Vera.ttf",
            Path(reportlab.__file__).resolve().parent / "fonts" / "VeraBd.ttf",
        ),
        (
            Path("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf"),
            Path("/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf"),
        ),
        (
            Path("C:/Windows/Fonts/arial.ttf"),
            Path("C:/Windows/Fonts/arialbd.ttf"),
        ),
    ]

    for regular, bold in candidates:
        if not regular.exists():
            continue
        try:
            if "ChatOmniRegular" not in pdfmetrics.getRegisteredFontNames():
                pdfmetrics.registerFont(TTFont("ChatOmniRegular", str(regular)))
            if bold.exists() and "ChatOmniBold" not in pdfmetrics.getRegisteredFontNames():
                pdfmetrics.registerFont(TTFont("ChatOmniBold", str(bold)))
            return "ChatOmniRegular", "ChatOmniBold" if bold.exists() else "ChatOmniRegular"
        except Exception:
            continue

    return "Helvetica", "Helvetica-Bold"


def _pdf_text(value) -> str:
    return html.escape(_text(value)).replace("\n", "<br/>")


def _write_pdf(path: Path, spec: dict):
    regular_font, bold_font = _pdf_fonts()
    base = getSampleStyleSheet()

    styles = {
        "title": ParagraphStyle(
            "ArtifactTitle", parent=base["Title"], fontName=bold_font,
            fontSize=20, leading=24, alignment=TA_CENTER, spaceAfter=10,
        ),
        "subtitle": ParagraphStyle(
            "ArtifactSubtitle", parent=base["Normal"], fontName=regular_font,
            fontSize=11, leading=15, alignment=TA_CENTER,
            textColor=colors.grey, spaceAfter=12,
        ),
        "heading": ParagraphStyle(
            "ArtifactHeading", parent=base["Heading2"], fontName=bold_font,
            fontSize=14, leading=18, spaceBefore=8, spaceAfter=5,
        ),
        "body": ParagraphStyle(
            "ArtifactBody", parent=base["BodyText"], fontName=regular_font,
            fontSize=10.5, leading=15, spaceAfter=4,
        ),
        "bullet": ParagraphStyle(
            "ArtifactBullet", parent=base["BodyText"], fontName=regular_font,
            fontSize=10.5, leading=15, leftIndent=12, firstLineIndent=-7,
        ),
        "table": ParagraphStyle(
            "ArtifactTable", parent=base["BodyText"], fontName=regular_font,
            fontSize=8.5, leading=11,
        ),
        "table_header": ParagraphStyle(
            "ArtifactTableHeader", parent=base["BodyText"], fontName=bold_font,
            fontSize=8.5, leading=11,
        ),
    }

    document = SimpleDocTemplate(
        str(path), pagesize=A4,
        leftMargin=20 * mm, rightMargin=20 * mm,
        topMargin=18 * mm, bottomMargin=18 * mm,
        title=_text(spec.get("title")), author="ChatOmni",
    )
    story = []
    available_width = A4[0] - 40 * mm

    title = _text(spec.get("title")).strip()
    subtitle = _text(spec.get("subtitle")).strip()
    if title:
        story.append(Paragraph(_pdf_text(title), styles["title"]))
    if subtitle:
        story.append(Paragraph(_pdf_text(subtitle), styles["subtitle"]))

    content = _text(spec.get("content")).strip()
    if content:
        for paragraph_text in re.split(r"\n\s*\n", content):
            if paragraph_text.strip():
                story.append(Paragraph(_pdf_text(paragraph_text.strip()), styles["body"]))
                story.append(Spacer(1, 2 * mm))

    for block in _document_blocks(spec):
        if not isinstance(block, dict):
            raise ValueError("Each section must be an object.")

        heading = _text(block.get("heading") or block.get("title")).strip()
        if heading:
            story.append(Paragraph(_pdf_text(heading), styles["heading"]))

        paragraphs = block.get("paragraphs") or []
        if isinstance(paragraphs, str):
            paragraphs = [paragraphs]
        for item in paragraphs:
            story.append(Paragraph(_pdf_text(item), styles["body"]))

        bullets = block.get("bullets") or []
        if isinstance(bullets, str):
            bullets = [bullets]
        for item in bullets:
            story.append(Paragraph(_pdf_text(item), styles["bullet"], bulletText="•"))

        numbered = block.get("numbered") or []
        if isinstance(numbered, str):
            numbered = [numbered]
        for index, item in enumerate(numbered, start=1):
            story.append(Paragraph(_pdf_text(item), styles["bullet"], bulletText=f"{index}."))

        for table_spec in block.get("tables") or []:
            headers, rows = _table_parts(table_spec)
            column_count = max([len(headers), *(len(row) for row in rows), 1])
            data = []

            if headers:
                data.append([
                    Paragraph(_pdf_text(headers[index] if index < len(headers) else ""), styles["table_header"])
                    for index in range(column_count)
                ])

            for row in rows:
                data.append([
                    Paragraph(_pdf_text(row[index] if index < len(row) else ""), styles["table"])
                    for index in range(column_count)
                ])

            if not data:
                data = [[Paragraph("", styles["table"])]]

            table = Table(
                data,
                colWidths=[available_width / column_count] * column_count,
                repeatRows=1 if headers else 0,
            )
            commands = [
                ("VALIGN", (0, 0), (-1, -1), "TOP"),
                ("GRID", (0, 0), (-1, -1), 0.5, colors.grey),
                ("LEFTPADDING", (0, 0), (-1, -1), 6),
                ("RIGHTPADDING", (0, 0), (-1, -1), 6),
                ("TOPPADDING", (0, 0), (-1, -1), 5),
                ("BOTTOMPADDING", (0, 0), (-1, -1), 5),
            ]
            if headers:
                commands.append(("BACKGROUND", (0, 0), (-1, 0), colors.whitesmoke))
            table.setStyle(TableStyle(commands))
            story.extend([table, Spacer(1, 4 * mm)])

    if not story:
        story.append(Paragraph("", styles["body"]))

    document.build(story)


def _write_xlsx(path: Path, spec: dict):
    sheets = spec.get("sheets")
    if sheets is None:
        sheets = [{
            "name": spec.get("sheet_name") or "Sheet1",
            "headers": spec.get("headers") or spec.get("columns") or [],
            "rows": spec.get("rows") or [],
        }]

    if not isinstance(sheets, list) or not sheets:
        raise ValueError("'sheets' must be a non-empty list.")

    workbook = Workbook()
    workbook.remove(workbook.active)
    used_names = set()

    for index, sheet_spec in enumerate(sheets, start=1):
        if not isinstance(sheet_spec, dict):
            raise ValueError("Each sheet must be an object.")

        requested_name = _text(sheet_spec.get("name") or f"Sheet{index}")
        name = re.sub(r"[\[\]:*?/\\]", "_", requested_name).strip()[:31] or f"Sheet{index}"
        base_name = name
        counter = 2
        while name in used_names:
            suffix = f"_{counter}"
            name = base_name[: 31 - len(suffix)] + suffix
            counter += 1
        used_names.add(name)

        worksheet = workbook.create_sheet(name)
        headers, rows = _table_parts(sheet_spec)

        if headers:
            for column, value in enumerate(headers, start=1):
                cell = worksheet.cell(1, column, value)
                cell.font = Font(bold=True)
                cell.alignment = Alignment(vertical="top", wrap_text=True)
            worksheet.freeze_panes = "A2"
            worksheet.auto_filter.ref = f"A1:{get_column_letter(len(headers))}1"

        start_row = 2 if headers else 1
        for row_index, row in enumerate(rows, start=start_row):
            for column, value in enumerate(row, start=1):
                if not isinstance(value, (str, int, float, bool)) and value is not None:
                    value = json.dumps(value, ensure_ascii=False)
                cell = worksheet.cell(row_index, column, value)
                cell.alignment = Alignment(vertical="top", wrap_text=True)

        max_columns = max([len(headers), *(len(row) for row in rows), 1])
        for column in range(1, max_columns + 1):
            longest = 0
            for row in range(1, min(worksheet.max_row, 200) + 1):
                value = worksheet.cell(row, column).value
                if value is not None:
                    longest = max(longest, len(str(value)))
            worksheet.column_dimensions[get_column_letter(column)].width = min(max(longest + 2, 10), 50)

    workbook.save(path)


def _write_pptx(path: Path, spec: dict):
    presentation = Presentation()
    presentation.core_properties.title = _text(spec.get("title"))
    presentation.core_properties.author = "ChatOmni"

    title = _text(spec.get("title")).strip()
    subtitle = _text(spec.get("subtitle")).strip()

    if title or subtitle:
        slide = presentation.slides.add_slide(presentation.slide_layouts[0])
        if slide.shapes.title:
            slide.shapes.title.text = title
        if len(slide.placeholders) > 1:
            slide.placeholders[1].text = subtitle

    slides = spec.get("slides") or []
    if not isinstance(slides, list):
        raise ValueError("'slides' must be a list.")

    for slide_spec in slides:
        if not isinstance(slide_spec, dict):
            raise ValueError("Each slide must be an object.")

        table_spec = slide_spec.get("table")
        layout = presentation.slide_layouts[5] if table_spec else presentation.slide_layouts[1]
        slide = presentation.slides.add_slide(layout)

        if slide.shapes.title:
            slide.shapes.title.text = _text(slide_spec.get("title"))

        if table_spec:
            headers, rows = _table_parts(table_spec)
            column_count = max([len(headers), *(len(row) for row in rows), 1])
            row_count = max(len(rows) + (1 if headers else 0), 1)
            table = slide.shapes.add_table(
                row_count, column_count,
                Inches(0.6), Inches(1.6), Inches(12.1), Inches(5.2),
            ).table

            row_cursor = 0
            if headers:
                for column in range(column_count):
                    cell = table.cell(0, column)
                    cell.text = _text(headers[column] if column < len(headers) else "")
                    for paragraph in cell.text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.font.bold = True
                row_cursor = 1

            for row in rows:
                for column in range(column_count):
                    table.cell(row_cursor, column).text = _text(row[column] if column < len(row) else "")
                row_cursor += 1
            continue

        body = None
        for placeholder in slide.placeholders:
            if placeholder != slide.shapes.title and hasattr(placeholder, "text_frame"):
                body = placeholder
                break
        if body is None:
            continue

        text_frame = body.text_frame
        text_frame.clear()

        items = []
        body_text = _text(slide_spec.get("body") or slide_spec.get("text")).strip()
        if body_text:
            items.append((body_text, 0))

        bullets = slide_spec.get("bullets") or []
        if isinstance(bullets, str):
            bullets = [bullets]
        for bullet in bullets:
            if isinstance(bullet, dict):
                items.append((_text(bullet.get("text")), int(bullet.get("level", 0))))
            else:
                items.append((_text(bullet), 0))

        for index, (text, level) in enumerate(items or [("", 0)]):
            paragraph = text_frame.paragraphs[0] if index == 0 else text_frame.add_paragraph()
            paragraph.text = text
            paragraph.level = max(0, min(level, 8))

    if len(presentation.slides) == 0:
        presentation.slides.add_slide(presentation.slide_layouts[6])

    presentation.save(path)


def _safe_zip_path(raw_path: str) -> str:
    if not isinstance(raw_path, str):
        raise ValueError("ZIP entry path must be a string.")

    path = PurePosixPath(raw_path.strip().replace("\\", "/"))
    if not raw_path.strip() or path.is_absolute() or ".." in path.parts or not path.name:
        raise ValueError("Unsafe or invalid path inside ZIP.")
    return path.as_posix()


def _write_zip(path: Path, spec: dict):
    files = spec.get("files") or []
    if isinstance(files, dict):
        files = [{"path": name, "content": content} for name, content in files.items()]

    if not isinstance(files, list) or not files:
        raise ValueError("ZIP spec must contain a non-empty 'files' list or object.")
    if len(files) > MAX_ZIP_FILES:
        raise ValueError("ZIP contains too many generated files.")

    seen = set()
    total_bytes = 0

    with zipfile.ZipFile(path, "w", compression=zipfile.ZIP_DEFLATED) as archive:
        for file_spec in files:
            if not isinstance(file_spec, dict):
                raise ValueError("Each ZIP entry must be an object.")

            archive_path = _safe_zip_path(file_spec.get("path") or file_spec.get("name") or "")
            if archive_path in seen:
                raise ValueError(f"Duplicate ZIP path: {archive_path}")
            seen.add(archive_path)

            content = _text(file_spec.get("content", ""))
            content_bytes = content.encode("utf-8")
            total_bytes += len(content_bytes)
            if total_bytes > MAX_ZIP_TEXT_BYTES:
                raise ValueError("Generated ZIP content is too large.")

            archive.writestr(archive_path, content_bytes)


def build_artifact(path: Path, suffix: str, spec: dict):
    if not isinstance(spec, dict):
        raise ValueError("Artifact spec must be an object.")

    handlers = {
        ".docx": _write_docx,
        ".pdf": _write_pdf,
        ".xlsx": _write_xlsx,
        ".pptx": _write_pptx,
        ".zip": _write_zip,
    }

    handler = handlers.get(suffix.lower())
    if handler is None:
        raise ValueError("Unsupported artifact extension.")

    handler(path, spec)